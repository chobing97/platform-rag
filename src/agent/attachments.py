"""첨부파일 전처리 — 이미지/PDF는 패스스루, Excel/PPT는 텍스트로 변환."""

import base64
import io
import logging
from dataclasses import dataclass

logger = logging.getLogger(__name__)

# MIME → 카테고리 매핑
_IMAGE_MIMES = {"image/png", "image/jpeg", "image/gif", "image/webp"}
_PDF_MIMES = {"application/pdf"}
_EXCEL_MIMES = {
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",  # .xlsx
    "application/vnd.ms-excel",  # .xls
}
_PPT_MIMES = {
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",  # .pptx
    "application/vnd.ms-powerpoint",  # .ppt
}

SUPPORTED_MIMES = _IMAGE_MIMES | _PDF_MIMES | _EXCEL_MIMES | _PPT_MIMES


@dataclass
class ContentBlock:
    """LLM에 전달할 콘텐츠 블록."""
    type: str  # "image", "document", "text"
    media_type: str | None = None  # MIME type (image, document)
    data: str | None = None  # base64 encoded (image, document)
    text: str | None = None  # 텍스트 내용 (text)
    file_name: str | None = None


def process_attachment(file_name: str, content_type: str, data: bytes) -> list[ContentBlock]:
    """첨부파일을 LLM 전달용 ContentBlock 리스트로 변환한다."""
    if content_type in _IMAGE_MIMES:
        return _passthrough_binary(file_name, content_type, data, block_type="image")

    if content_type in _PDF_MIMES:
        return _passthrough_binary(file_name, content_type, data, block_type="document")

    if content_type in _EXCEL_MIMES:
        return _convert_excel(file_name, data)

    if content_type in _PPT_MIMES:
        return _convert_pptx(file_name, data)

    logger.warning("지원하지 않는 파일 타입: %s (%s)", file_name, content_type)
    return [ContentBlock(type="text", text=f"[지원하지 않는 파일 형식: {file_name} ({content_type})]")]


def _passthrough_binary(file_name: str, content_type: str, data: bytes, *, block_type: str) -> list[ContentBlock]:
    """이미지/PDF를 base64로 인코딩하여 그대로 전달한다."""
    b64 = base64.standard_b64encode(data).decode("ascii")
    return [ContentBlock(type=block_type, media_type=content_type, data=b64, file_name=file_name)]


def _convert_excel(file_name: str, data: bytes) -> list[ContentBlock]:
    """Excel 파일을 마크다운 테이블로 변환한다."""
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts: list[str] = [f"**[첨부: {file_name}]**\n"]

    for sheet in wb.worksheets:
        rows = list(sheet.iter_rows(values_only=True))
        if not rows:
            continue

        parts.append(f"### 시트: {sheet.title}\n")

        # 헤더
        header = [str(c) if c is not None else "" for c in rows[0]]
        parts.append("| " + " | ".join(header) + " |")
        parts.append("| " + " | ".join(["---"] * len(header)) + " |")

        # 데이터 (최대 200행 — LLM 컨텍스트 절약)
        for row in rows[1:201]:
            cells = [str(c) if c is not None else "" for c in row]
            # 열 수 맞춤
            while len(cells) < len(header):
                cells.append("")
            parts.append("| " + " | ".join(cells[:len(header)]) + " |")

        if len(rows) > 201:
            parts.append(f"\n> ... 총 {len(rows) - 1}행 중 처음 200행만 표시\n")
        parts.append("")

    wb.close()
    return [ContentBlock(type="text", text="\n".join(parts), file_name=file_name)]


def _convert_pptx(file_name: str, data: bytes) -> list[ContentBlock]:
    """PPT 파일에서 슬라이드별 텍스트와 이미지를 추출한다."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(io.BytesIO(data))
    blocks: list[ContentBlock] = []
    text_parts: list[str] = [f"**[첨부: {file_name}]**\n"]

    for i, slide in enumerate(prs.slides, 1):
        slide_texts: list[str] = []
        for shape in slide.shapes:
            # 텍스트 추출
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        slide_texts.append(t)

            # 테이블 추출
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    slide_texts.append("| " + " | ".join(cells) + " |")

            # 이미지 추출
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img = shape.image
                mime = img.content_type
                if mime in _IMAGE_MIMES:
                    b64 = base64.standard_b64encode(img.blob).decode("ascii")
                    blocks.append(ContentBlock(
                        type="image", media_type=mime, data=b64,
                        file_name=f"{file_name}_slide{i}_img",
                    ))

        if slide_texts:
            text_parts.append(f"### 슬라이드 {i}")
            text_parts.extend(slide_texts)
            text_parts.append("")

    # 텍스트 블록을 맨 앞에 배치
    blocks.insert(0, ContentBlock(type="text", text="\n".join(text_parts), file_name=file_name))
    return blocks
